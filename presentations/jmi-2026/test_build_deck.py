"""Test suite for deck builder (Task 11)."""
import os
from pathlib import Path
from build_deck import build


def test_build_produces_pptx_with_correct_slide_count(tmp_path):
    out = tmp_path / "deck.pptx"
    n = build(out_path=out)
    assert out.exists()
    # Slide count is data-driven: number of entries in copy.yaml.
    import yaml
    spec = yaml.safe_load(Path("copy.yaml").read_text())
    assert n == len(spec["slides"])


def test_no_underscore_or_caret_anywhere_in_deck(tmp_path):
    from pptx import Presentation
    out = tmp_path / "deck.pptx"
    build(out_path=out)
    prs = Presentation(str(out))
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    assert "_" not in run.text, f"slide {i}: underscore leak: {run.text!r}"
                    assert "^" not in run.text, f"slide {i}: caret leak: {run.text!r}"

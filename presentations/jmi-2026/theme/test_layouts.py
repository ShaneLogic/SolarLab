"""Layout builder tests."""
from pptx import Presentation
from theme.layouts import add_cover, add_divider, add_content


def test_cover_has_title_and_two_subline_runs():
    prs = Presentation()
    prs.slide_width = 9144000          # 10in
    prs.slide_height = 5143500         # 5.625in (16:9)
    add_cover(prs, title="Photovoltaic material screening",
                  authors="Xuan-Yan Chen, ... (HKUST-GZ)",
                  affiliation="JMI 2026")
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    texts = [shape.text_frame.text for shape in slide.shapes if shape.has_text_frame]
    joined = " | ".join(texts)
    assert "Photovoltaic material screening" in joined
    assert "Xuan-Yan Chen" in joined
    assert "JMI 2026" in joined


def test_divider_writes_numeral_and_title():
    prs = Presentation()
    prs.slide_width, prs.slide_height = 9144000, 5143500
    add_divider(prs, number="01", title="PROBLEM")
    slide = prs.slides[0]
    joined = " | ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "01" in joined
    assert "PROBLEM" in joined


def test_content_renders_subscripts_via_runs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = 9144000, 5143500
    add_content(prs, section_label="03 · ML & FEATURES", slide_index=10,
                total=22, title="Validation",
                subtitle="Predicted vs HSE06",
                bullets=["R[sup]2[/sup] (E[sub]g[/sub]) = 0.91"])
    slide = prs.slides[0]
    # Walk every run in every textbox; ensure no '_' or '^' literal survived
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                assert "_" not in run.text, f"underscore leaked: {run.text!r}"
                assert "^" not in run.text, f"caret leaked: {run.text!r}"

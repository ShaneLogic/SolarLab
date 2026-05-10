"""Cover, divider, content layout builders (spec § 7.4 & § 7.5)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from theme import palette, fonts
from theme.runs import parse_segments, write_segments


_BLANK_LAYOUT = 6        # PowerPoint blank layout index


def _add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[_BLANK_LAYOUT])


def _add_textbox(slide, left, top, width, height, text, role,
                 color, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    write_segments(p, parse_segments(text), role, color)
    return tb


def add_cover(prs, *, title, authors, affiliation):
    slide = _add_blank(prs)
    sw, sh = prs.slide_width, prs.slide_height
    _add_textbox(slide, Emu(0), Emu(int(sh * 0.40)), sw, Inches(1),
                 title, fonts.TITLE, palette.INK, PP_ALIGN.CENTER)
    _add_textbox(slide, Emu(0), Emu(int(sh * 0.55)), sw, Inches(0.5),
                 authors, fonts.SUBTITLE, palette.BODY, PP_ALIGN.CENTER)
    _add_textbox(slide, Emu(0), Emu(int(sh * 0.62)), sw, Inches(0.4),
                 affiliation, fonts.CAPTION, palette.HAIRLINE, PP_ALIGN.CENTER)
    # accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Emu(int(sw * 0.42)), Emu(int(sh * 0.52)),
                                  Emu(int(sw * 0.16)), Emu(int(sh * 0.004)))
    line.fill.solid()
    line.fill.fore_color.rgb = palette.ACCENT
    line.line.fill.background()
    return slide


def add_divider(prs, *, number, title):
    """Full-bleed numeral divider (spec § 7.5)."""
    slide = _add_blank(prs)
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
    bg.fill.solid()
    bg.fill.fore_color.rgb = palette.INK
    bg.line.fill.background()
    _add_textbox(slide, Emu(int(sw * 0.06)), Emu(int(sh * 0.10)),
                 Emu(int(sw * 0.6)), Inches(4),
                 number, fonts.DIVIDER_NUMERAL, palette.ACCENT)
    _add_textbox(slide, Emu(int(sw * 0.06)), Emu(int(sh * 0.78)),
                 Emu(int(sw * 0.88)), Inches(0.6),
                 title, fonts.DIVIDER_TITLE, palette.BODY)
    return slide


def add_content(prs, *, section_label, slide_index, total, title,
                subtitle, bullets, figure_path=None):
    slide = _add_blank(prs)
    sw, sh = prs.slide_width, prs.slide_height
    # Layout band budget (inches, 16:9 deck = 10 × 5.625):
    #   header  : 0.30 .. 0.60  (height 0.30)
    #   title   : 0.70 .. 1.40  (height 0.70 — accommodates 2 lines at 22pt)
    #   subtitle: 1.45 .. 1.90  (height 0.45)
    #   body    : 2.00 .. 5.00  (height 3.00 — figure / bullets)
    #   footer  : 5.10 .. 5.40  (height 0.30 — brand mark)
    # All bands are non-overlapping by 0.05 in margins.

    # header row
    _add_textbox(slide, Inches(0.5), Inches(0.3), Inches(7), Inches(0.3),
                 section_label.upper(), fonts.SECTION_LABEL, palette.BODY)
    _add_textbox(slide, Inches(7.5), Inches(0.3), Inches(2), Inches(0.3),
                 f"{slide_index} / {total}", fonts.SLIDE_COUNTER,
                 palette.HAIRLINE, align=PP_ALIGN.RIGHT)
    # body title + subtitle
    _add_textbox(slide, Inches(0.5), Inches(0.7), Inches(9), Inches(0.7),
                 title, fonts.CONTENT_TITLE, palette.INK)
    _add_textbox(slide, Inches(0.5), Inches(1.45), Inches(9), Inches(0.45),
                 subtitle, fonts.SUBTITLE, palette.BODY)
    # figure (left) + bullets (right) — or full-width figure if no bullets
    body_top_in = 2.0
    body_h_in = 3.0
    has_bullets = bool(bullets)
    if figure_path:
        from PIL import Image
        with Image.open(str(figure_path)) as im:
            src_w, src_h = im.size
        # Figure-only slides (no bullets) use the full 9-inch body width;
        # otherwise the figure occupies the 5.4-inch left column.
        if has_bullets:
            max_w_in, max_h_in = 5.4, body_h_in
            box_left_in = 0.5
        else:
            max_w_in, max_h_in = 9.0, body_h_in
            box_left_in = 0.5
        src_aspect = src_w / src_h
        if src_aspect > (max_w_in / max_h_in):
            fig_w_in = max_w_in
            fig_h_in = max_w_in / src_aspect
        else:
            fig_h_in = max_h_in
            fig_w_in = max_h_in * src_aspect
        x_in = box_left_in + (max_w_in - fig_w_in) / 2.0
        y_in = body_top_in + (max_h_in - fig_h_in) / 2.0
        slide.shapes.add_picture(str(figure_path),
                                 Inches(x_in), Inches(y_in),
                                 width=Inches(fig_w_in),
                                 height=Inches(fig_h_in))
        bullet_left = Inches(6.2)
        bullet_w = Inches(3.3)
    else:
        bullet_left = Inches(0.5)
        bullet_w = Inches(9)
    if bullets:
        tb = slide.shapes.add_textbox(bullet_left, Inches(body_top_in),
                                      bullet_w, Inches(body_h_in))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(2)
            write_segments(p, parse_segments(f"•  {line}"),
                           fonts.BODY, palette.BODY)
    # footer — brand mark only (author cite removed per user request)
    _add_textbox(slide, Inches(7.5), Inches(5.10), Inches(2), Inches(0.30),
                 "SOLARLAB", fonts.FOOTER, palette.ACCENT, align=PP_ALIGN.RIGHT)
    return slide

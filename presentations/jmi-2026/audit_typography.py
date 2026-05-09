"""Walk every run of output/jmi-2026.pptx and audit sub/sup runs."""
from pathlib import Path
from pptx import Presentation
from pptx.oxml.ns import qn

DECK = Path("output/jmi-2026.pptx")


def main():
    prs = Presentation(str(DECK))
    runs_total = 0
    runs_sub = 0
    runs_sup = 0
    leaks = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    runs_total += 1
                    if "_" in run.text or "^" in run.text:
                        leaks.append((slide_idx, run.text))
                    rpr = run._r.find(qn("a:rPr"))
                    if rpr is not None:
                        bl = rpr.get("baseline")
                        if bl == "-25000":
                            runs_sub += 1
                        elif bl == "30000":
                            runs_sup += 1
    print(f"total runs: {runs_total}")
    print(f"  subscript runs: {runs_sub}")
    print(f"  superscript runs: {runs_sup}")
    if leaks:
        for s, t in leaks:
            print(f"  LEAK on slide {s}: {t!r}")
        raise SystemExit(1)
    print("PASS: no _ or ^ literals; sub/sup runs have baseline attr.")


if __name__ == "__main__":
    main()

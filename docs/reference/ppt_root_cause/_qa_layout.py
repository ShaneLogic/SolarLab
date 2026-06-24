#!/usr/bin/env python3
"""Programmatic layout + typography QA for the SCAPS root-cause deck."""
import os, re
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

PPTX = "/Users/shane/Library/CloudStorage/OneDrive-HKUST(Guangzhou)/SolarLab/docs/reference/ppt_root_cause/SolarLab_SCAPS_root_cause.pptx"
FIGDIR = "/Users/shane/Library/CloudStorage/OneDrive-HKUST(Guangzhou)/SolarLab/docs/reference/ppt_root_cause/figures"

EMU_PER_IN = 914400
prs = Presentation(PPTX)
SW, SH = prs.slide_width, prs.slide_height
print(f"SLIDE: {SW/EMU_PER_IN:.3f} x {SH/EMU_PER_IN:.3f} in ; {len(prs.slides)} slides")
MARGIN = int(0.5 * EMU_PER_IN)

# physics tokens that should be formatted (sub/superscript) if they appear as literal runs
TOKEN_RE = re.compile(r'(_\{|\^\{|\bVoc\b|\bJsc\b|\bJ0\b|\bV_bi\b|cm-3|cm\^-3|m\^2|delta-E_C|Nd_ETL|Nt_C|Nt_V|N_grid|E_C|E_V)')
NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

def rect(sp):
    try:
        return (sp.left, sp.top, sp.width, sp.height)
    except Exception:
        return None

def overlap_area(a, b):
    if a is None or b is None: return 0
    al, at, aw, ah = a; bl, bt, bw, bh = b
    if None in a or None in b: return 0
    ix = max(0, min(al+aw, bl+bw) - max(al, bl))
    iy = max(0, min(at+ah, bt+bh) - max(at, bt))
    return ix * iy

def shp_label(sp):
    nm = sp.name
    ph = ""
    try:
        if sp.is_placeholder:
            ph = f"/ph={sp.placeholder_format.type}"
    except Exception:
        pass
    kind = "PIC" if sp.shape_type == MSO_SHAPE_TYPE.PICTURE else ("TXT" if sp.has_text_frame and sp.text_frame.text.strip() else "SHP")
    return f"{kind}:{nm}{ph}"

# source PNG aspect ratios
png_aspect = {}
for f in os.listdir(FIGDIR):
    if f.endswith(".png"):
        try:
            with Image.open(os.path.join(FIGDIR, f)) as im:
                png_aspect[f] = im.size[0]/im.size[1]
        except Exception: pass

issues = []
TOL = int(0.04*EMU_PER_IN)**2  # ~tiny overlap tolerance in area (0.04in x 0.04in)
LINTOL = int(0.03*EMU_PER_IN)

for si, slide in enumerate(prs.slides, 1):
    shapes = [s for s in slide.shapes]
    geoms = [(s, rect(s)) for s in shapes]
    # title detection
    title_idx = set()
    for i, s in enumerate(shapes):
        try:
            if s.is_placeholder and 'TITLE' in str(s.placeholder_format.type):
                title_idx.add(i)
        except Exception: pass

    # (a) OVERLAP — pairwise, focus text-vs-pic / pic-vs-title
    for i in range(len(geoms)):
        for j in range(i+1, len(geoms)):
            (sa, ga), (sb, gb) = geoms[i], geoms[j]
            ar = overlap_area(ga, gb)
            if ar > TOL:
                a_pic = sa.shape_type == MSO_SHAPE_TYPE.PICTURE
                b_pic = sb.shape_type == MSO_SHAPE_TYPE.PICTURE
                a_txt = sa.has_text_frame and sa.text_frame.text.strip()
                b_txt = sb.has_text_frame and sb.text_frame.text.strip()
                interesting = (a_pic and b_txt) or (b_pic and a_txt) or (a_pic and b_pic) or (i in title_idx) or (j in title_idx)
                # only report meaningful overlaps (skip group-member nesting where one fully contains other and is a group/connector backdrop)
                sev = "high" if interesting else "med"
                in2 = ar/(EMU_PER_IN**2)
                if in2 > 0.02:
                    issues.append((si, sev, "OVERLAP",
                        f"{shp_label(sa)} <> {shp_label(sb)} overlap {in2:.3f} in^2"))

    # (b) OFF-SLIDE / MARGIN
    for s, g in geoms:
        if g is None or None in g: continue
        l, t, w, h = g
        if l < -LINTOL or t < -LINTOL or l+w > SW+LINTOL or t+h > SH+LINTOL:
            issues.append((si, "high", "OFF-SLIDE",
                f"{shp_label(s)} bbox [{l/EMU_PER_IN:.2f},{t/EMU_PER_IN:.2f},{(l+w)/EMU_PER_IN:.2f},{(t+h)/EMU_PER_IN:.2f}] exceeds slide"))
        elif l < MARGIN-LINTOL or t < MARGIN-LINTOL or l+w > SW-MARGIN+LINTOL or t+h > SH-MARGIN+LINTOL:
            # margin breach only flagged low unless it's far
            over = max(MARGIN-l, MARGIN-t, (l+w)-(SW-MARGIN), (t+h)-(SH-MARGIN))/EMU_PER_IN
            if over > 0.08:
                issues.append((si, "low", "MARGIN",
                    f"{shp_label(s)} breaches 0.5in safe margin by {over:.2f}in"))

    # (c)+(d) TEXT scan
    is_ref_slide = False
    for s in shapes:
        if s.has_text_frame:
            tt = s.text_frame.text.lower()
            if 'reference' in tt and si >= len(prs.slides)-1:
                is_ref_slide = True

    for s in shapes:
        if not s.has_text_frame: continue
        tf = s.text_frame
        # autosize shrink detection
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
            if tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
                issues.append((si, "low", "AUTOSIZE", f"{shp_label(s)} text frame uses shrink-to-fit autosize"))
        except Exception: pass
        for p in tf.paragraphs:
            for r in p.runs:
                txt = r.text
                if not txt.strip(): continue
                # literal token check
                m = TOKEN_RE.search(txt)
                if m:
                    # check if this run carries baseline (sub/sup) — if so it's intended
                    rpr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    has_base = rpr is not None and rpr.get('baseline') is not None
                    if not has_base:
                        issues.append((si, "med", "TYPO",
                            f"unformatted token '{m.group(0)}' in run: \"{txt[:50]}\""))
                # font size check
                sz = r.font.size
                if sz is None:
                    # inherit from paragraph / placeholder — check pPr defRPr
                    pass
                else:
                    pt = sz.pt
                    if pt < 12 and not is_ref_slide:
                        issues.append((si, "med", "FONT",
                            f"{pt:.1f}pt (<12) run: \"{txt[:40]}\""))

    # (d) picture aspect-ratio distortion
    for s in shapes:
        if s.shape_type != MSO_SHAPE_TYPE.PICTURE: continue
        try:
            disp = s.width / s.height
        except Exception:
            continue
        # match source png by image blob filename
        src = None
        try:
            img = s.image
            src = os.path.basename(img.filename) if img.filename else None
        except Exception:
            src = None
        srca = png_aspect.get(src) if src else None
        if srca:
            rel = abs(disp - srca)/srca
            if rel > 0.03:
                issues.append((si, "med" if rel>0.06 else "low", "DISTORT",
                    f"{src} displayed AR {disp:.3f} vs source {srca:.3f} ({rel*100:.1f}% off)"))

# verify intended sub/sup runs DO carry baseline (count them so we know formatting was applied somewhere)
base_runs = 0
for slide in prs.slides:
    for s in slide.shapes:
        if not s.has_text_frame: continue
        for p in s.text_frame.paragraphs:
            for r in p.runs:
                rpr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                if rpr is not None and rpr.get('baseline') is not None:
                    base_runs += 1
print(f"Runs carrying baseline (sub/sup) attribute: {base_runs}")

# dump per-slide shape inventory (compact)
print("\n--- SHAPE INVENTORY ---")
for si, slide in enumerate(prs.slides, 1):
    n_pic = sum(1 for s in slide.shapes if s.shape_type==MSO_SHAPE_TYPE.PICTURE)
    n_txt = sum(1 for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip())
    title = ""
    for s in slide.shapes:
        try:
            if s.is_placeholder and 'TITLE' in str(s.placeholder_format.type):
                title = s.text_frame.text[:55]
        except Exception: pass
    print(f"  S{si}: {len(slide.shapes)} shapes ({n_pic} pic, {n_txt} txt) | {title}")

print("\n--- ISSUES ---")
order = {"high":0,"med":1,"low":2}
for it in sorted(issues, key=lambda x:(x[0], order[x[1]])):
    print(f"S{it[0]} [{it[1].upper()}] {it[2]}: {it[3]}")
print(f"\nTOTAL ISSUES: {len(issues)}")

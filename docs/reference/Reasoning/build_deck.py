#!/usr/bin/env python3
"""Build the SolarLab vs SCAPS-1D root-cause analysis deck (16:9)."""
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")
OUT = os.path.join(HERE, "SolarLab_SCAPS_root_cause.pptx")
GEN_DATE = "2026-06-04"

# Palette
NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLUE = RGBColor(0x1A, 0x4E, 0x8A)
SLATE = RGBColor(0x33, 0x3A, 0x4A)
GREEN = RGBColor(0x1B, 0x7A, 0x3D)
RED = RGBColor(0xB0, 0x2A, 0x2A)
GREY = RGBColor(0x5A, 0x60, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- rich-text markup helper ----------
_TOKEN = re.compile(r"(_\{[^}]*\}|\^\{[^}]*\})")


def _set_baseline(run, val):
    rPr = run._r.get_or_add_rPr()
    rPr.set("baseline", str(val))


def _add_run(p, text, size, bold, color, baseline=None):
    run = p.add_run()
    run.text = text
    f = run.font
    if baseline is not None:
        f.size = Pt(size * 0.72)
        _set_baseline(run, baseline)
    else:
        f.size = Pt(size)
    f.bold = bold
    if color is not None:
        f.color.rgb = color
    return run


def add_rich(p, text, size=18, bold=False, color=None):
    """Parse _{...}/^{...} markup into runs with real sub/superscript baseline."""
    pos = 0
    for m in _TOKEN.finditer(text):
        if m.start() > pos:
            _add_run(p, text[pos:m.start()], size, bold, color)
        seg = m.group(0)
        inner = seg[2:-1]
        if seg.startswith("_{"):
            _add_run(p, inner, size, bold, color, baseline=-25000)
        else:
            _add_run(p, inner, size, bold, color, baseline=30000)
        pos = m.end()
    if pos < len(text):
        _add_run(p, text[pos:], size, bold, color)


# ---------- layout primitives ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def add_title(slide, text, size=30):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    add_rich(p, text, size=size, bold=True, color=NAVY)
    # accent underline bar
    bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.32), Inches(12.1), Pt(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    return tb


def add_bullets(slide, bullets, left, top, width, height, size=17):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(0)
        p.line_spacing = 1.05
        add_rich(p, "•  ", size=size, bold=False, color=BLUE)
        add_rich(p, b, size=size, bold=False, color=SLATE)
    return tb


def add_picture_fit(slide, path, left, top, box_w, box_h, missing_note):
    if not os.path.exists(path):
        missing_note.append(os.path.basename(path))
        box = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(box_w), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
        box.line.color.rgb = RED
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        add_rich(p, "MISSING FIGURE\n" + os.path.basename(path), size=14, bold=True, color=RED)
        return
    w, h = Image.open(path).size
    aspect = w / h
    draw_w = box_w
    draw_h = draw_w / aspect
    if draw_h > box_h:
        draw_h = box_h
        draw_w = draw_h * aspect
    off_l = left + (box_w - draw_w) / 2.0
    off_t = top + (box_h - draw_h) / 2.0
    slide.shapes.add_picture(path, Inches(off_l), Inches(off_t), width=Inches(draw_w), height=Inches(draw_h))


def add_notes(slide, notes):
    slide.notes_slide.notes_text_frame.text = notes


# Figure-slide standard geometry
BUL_L, BUL_T, BUL_W, BUL_H = 0.6, 1.55, 5.5, 5.4
PIC_L, PIC_T, PIC_W, PIC_H = 6.35, 1.55, 6.4, 5.4
TXT_L, TXT_T, TXT_W, TXT_H = 0.8, 1.6, 11.7, 5.3

missing = []


def figure_slide(title, fig, bullets, notes, bsize=16):
    s = new_slide()
    add_title(s, title)
    add_bullets(s, bullets, BUL_L, BUL_T, BUL_W, BUL_H, size=bsize)
    add_picture_fit(s, os.path.join(FIG, fig), PIC_L, PIC_T, PIC_W, PIC_H, missing)
    add_notes(s, notes)
    return s


def text_slide(title, bullets, notes, bsize=19):
    s = new_slide()
    add_title(s, title)
    add_bullets(s, bullets, TXT_L, TXT_T, TXT_W, TXT_H, size=bsize)
    add_notes(s, notes)
    return s


# ===================== SLIDE 1 (title) =====================
s = new_slide()
# big title
tb = s.shapes.add_textbox(Inches(0.6), Inches(1.9), Inches(12.1), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.LEFT
add_rich(p, "SolarLab vs. SCAPS-1D", size=40, bold=True, color=NAVY)
p2 = tf.add_paragraph()
p2.space_before = Pt(4)
add_rich(p2, "Root-Cause Analysis of the Base-Model V_{oc} Discrepancy", size=27, bold=True, color=BLUE)
# subtitle line
sb = s.shapes.add_textbox(Inches(0.6), Inches(4.05), Inches(12.1), Inches(2.2))
stf = sb.text_frame
stf.word_wrap = True
subs = [
    "Drift-diffusion cross-validation on the reference MAPbI_{3} base stack",
    "Primary discrepancy: V_{oc} 1.072 V vs. 1.1676 V (−96 mV)",
    "Multi-agent adversarial review; trend-first evaluation standard",
    "One survivor of seven candidate root causes",
]
for i, b in enumerate(subs):
    p = stf.paragraphs[0] if i == 0 else stf.add_paragraph()
    p.space_after = Pt(6)
    add_rich(p, "•  ", size=18, color=BLUE)
    add_rich(p, b, size=18, color=SLATE)
# footer
fb = s.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12.1), Inches(0.4))
fp = fb.text_frame.paragraphs[0]
add_rich(fp, "Generated " + GEN_DATE + "  ·  SolarLab perovskite_sim cross-validation", size=12, color=GREY)
# accent bar top
bar = s.shapes.add_shape(1, Inches(0.6), Inches(1.7), Inches(4.2), Pt(4))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
add_notes(s, "This deck reports a cross-validation of two 1D drift-diffusion solvers, not a comparison against experiment. Neither code is \"truth\"; we anchor realism against the literature. The headline question is why SolarLab's Voc sits 96 mV below SCAPS, and what is actually fixable.")

# ===================== SLIDE 2 =====================
figure_slide(
    "Context: Two DD Solvers on the Reference MAPbI_{3} Base Model",
    "f4_solver_arch.png",
    [
        "Stack: glass / spiro HTL 20 nm / MAPbI_{3} 800 nm (E_{g}=1.53 eV) / TiO_{2} ETL 25 nm",
        "Conduction-band offset Delta-E_{C} = −0.16 eV (a \"cliff\")",
        "SCAPS: Gummel+Newton, (psi, E_{Fn}, E_{Fp}) [1]; SolarLab: method-of-lines Radau",
        "Both Boltzmann, Scharfetter–Gummel flux [5], tunnelling OFF [2]",
        "Same physics class; differ in boundary and interface treatments",
    ],
    "SCAPS-1D (Gent) is the de-facto reference; SolarLab is our Python perovskite_sim library. Both solve Poisson plus two continuity equations with SG flux. The divergences live in boundary conditions, interface recombination, and optics, not in the core formulation.",
)

# ===================== SLIDE 3 =====================
text_slide(
    "Method: Multi-Agent Adversarial Review, Trend-First Standard",
    [
        "Three lenses: physics-correctness, numerical-algorithm, data-consistency",
        "A cause must survive a majority of lenses to \"hold\"",
        "Parity judged by trend (direction + shape) first, absolutes second",
        "Closure = SolarLab V_{oc} span ÷ SCAPS V_{oc} span per sweep",
        "Every claim traced to source data or a cited reference",
    ],
    "Each candidate root cause was attacked from three independent angles. We deliberately avoid chasing absolute parity; the house rule prizes sweep direction and shape fidelity. This is why a 96 mV absolute gap can be acceptable while a single direction reversal is the real finding.",
)

# ===================== SLIDE 4 =====================
figure_slide(
    "Base-Model Figures of Merit: V_{oc}, J_{sc}, FF, PCE",
    "f5_fom_bars.png",
    [
        "V_{oc}: 1.072 vs. 1.1676 V (−96 mV, −8.2%) — primary gap",
        "J_{sc} genuine: 23.96 vs. 26.282 mA/cm^{2} (table 25.73 is foreign)",
        "FF: ~86.2% vs. 86.99% (−1.4 pp)",
        "PCE genuine: 22.1% vs. 26.69% (the 25.73 splice inflates the table)",
        "Self-consistent SolarLab: 1.0709·23.957·0.862 = 22.1%",
    ],
    "The quoted Jsc=25.73 is a contaminant spliced from a separate N_grid=30 run; the genuine self-consistent value is 23.96. Using the real Jsc, all four FOMs lock algebraically. SCAPS obeys PCE = Voc*Jsc*FF to within rounding; so does the genuine SolarLab tuple.",
)

# ===================== SLIDE 5 =====================
figure_slide(
    "The Ten-Sweep Landscape: Which Levers Diverge",
    "f6_survival.png",
    [
        "Ten parameter sweeps inherited from the SCAPS reference",
        "Matched: CHI_{ETL} (CBO) 83% span, correct direction",
        "Direction reversal: Nd_{ETL} donor doping — the key anomaly",
        "Dead: bulk N_{t} (Nt_{C}/Nt_{V}_PVK), iface trap energy",
        "Over-sensitive: Nt_{HTL}_PVK, bulk E_{t} (~8x)",
        "Span-ratio \"closure\" can mask a shape mismatch (Nt_{PVK}_ETL)",
    ],
    "The CBO sweep is SCAPS's dominant Voc lever and SolarLab tracks it well. Most divergences trace to the interface-SRH model or carrier-density suppression. The one qualitatively wrong result, a direction reversal, is the donor-doping sweep, and that is where the robust root cause lives.",
)

# ===================== SLIDE 6 =====================
figure_slide(
    "Adversarial Survival: One of Seven Root Causes Holds",
    "f6_survival.png",
    [
        "Nd_{ETL} donor-doping reversal: 3/3 lenses — THE SURVIVOR",
        "V_{oc} −96 mV deficit: 1/3 (mechanism contested)",
        "Nt_{C}/Nt_{V} bulk insensitivity: 1/3; CHI_{ETL} CBO: 1/3",
        "PCE/FF self-consistency: 1/3 (reporting artifact)",
        "Nt_{PVK}_ETL interface channel: 0/3; J_{sc} −2% optical: 0/3",
    ],
    "Only the donor-doping direction reversal survived all three lenses. The 96 mV deficit is real but its mechanism could not be pinned. Several headline claims, including the prior 37x J0 story and the \"-2% optical\" framing, did not survive at all.",
)

# ===================== SLIDE 7 =====================
figure_slide(
    "Root Cause: Frozen V_{bi} Severs the ETL-Doping Lever",
    "f1_voc_lever.png",
    [
        "Poisson boundary hard-wired to V_{bi}=1.30 V (config line 40)",
        "Band-derived V_{bi,eff} computed but routed only to sweep range",
        "Raising N_{D} cannot lift built-in potential → V_{oc} flat",
        "SCAPS rises +100 mV; SolarLab nets −11 mV (direction reversal)",
        "Signature: V_{oc} falls 1.0954→1.0656 V then drifts to 1.0844 V",
        "Classification: numerical boundary, not interface recombination",
    ],
    "The donor-doping lever is severed because the built-in potential is frozen at the boundary. The only residual N_D channel is the ohmic pin p_R = ni^2/N_D, which actually suppresses splitting, giving the wrong sign. Projection ON keeps the wrong sign, confirming it is not an interface-recombination effect.",
)

# ===================== SLIDE 8 =====================
figure_slide(
    "Refuting the Prior 37x-J_{0} Claim: The CBO Gap Persists",
    "f2_cbo_refute.png",
    [
        "Prior report: 37x J_{0} from QFL dissipated across band offsets",
        "Test: sweep the conduction-band offset to Delta-E_{C} → 0",
        "Removing the offset does NOT close the gap",
        "At Delta-E_{C}=0: SolarLab ~1.085 V vs. SCAPS ~1.248 V (~163 mV)",
        "Gap widens, not closes — the offset is not the seat of the deficit",
        "Prior central physical claim should be retracted",
    ],
    "If the deficit were driven by QFL loss across the band offset, removing the offset would close it. Instead the gap grows to ~163 mV at flat-band. SolarLab saturates near 1.085 V while SCAPS keeps climbing. This falsifies the prior mechanism and redirects attention to the interface-SRH model difference.",
)

# ===================== SLIDE 9 =====================
figure_slide(
    "The −96 mV V_{oc} Deficit: Physical but Contested Mechanism",
    "f7_experiment_voc.png",
    [
        "Deficit is real and physical (FF healthy ~0.878, not an artifact)",
        "Effective J_{0} ratio ~40x → V_{T}·ln(41.7) = 96.5 mV",
        "Contact-sink attribution REFUTED: same S→infinity limit as SCAPS",
        "Better candidate: SolarLab vs. SCAPS Pauwels–Vanhoutte interface SRH [7,8]",
        "Exonerated as drivers: V_{bi}, n_i, TE flux cap, FOM extraction",
        "Confidence: medium interface-SRH; low contact-BC specifically",
    ],
    "The symptom, a ~40x effective J0 mapping to ~96 mV, is robust. But the contact-sink mechanism fails the physics and data lenses; the gap is monotone in CBO and moves under interface-projection toggles. The most likely origin is the interface-SRH-model difference, which was not independently confirmed.",
)

# ===================== SLIDE 10 =====================
figure_slide(
    "PCE Anomaly = A J_{sc} Cell-Splice, Not a Solver Bug",
    "f3_pce_splice.png",
    [
        "PCE ≡ V_{oc}·J_{sc}·FF is algebraically locked per run",
        "Task tuple violates the lock by +1.51 pp → multi-source pairing",
        "Foreign cell is J_{sc}=25.73 (an N_{grid}=30 run), not PCE",
        "Genuine J_{sc}=23.96 → locked PCE = 22.1%",
        "Prior \"true PCE ~23.6%\" re-commits the splice — also wrong",
        "Classification: reporting artifact (held 3/3 at top level)",
    ],
    "The four FOMs must satisfy PCE = Voc*Jsc*FF in any single run. The reference table glued a coarse-grid Jsc onto a fine-grid Voc/FF/PCE triple. The fix is to re-extract all four from one J-V array; this alone erases the headline -4.6 pp PCE gap.",
)

# ===================== SLIDE 11 =====================
figure_slide(
    "Two Engines: Where SolarLab and SCAPS Diverge",
    "f4_solver_arch.png",
    [
        "V_{bi}: frozen scalar (SolarLab) vs. doping/band-tracking flatband (SCAPS) [1,6]",
        "Interface SRH: bulk-projected cross-carrier vs. single-plane Pauwels–Vanhoutte [8]",
        "Contacts: hard-Dirichlet ohmic vs. flatband work-function [6]",
        "Transport: TE flux-cap on SG vs. TE as the interface BC [4]",
        "Optics: TMM front reflection + finite absorber vs. R=0 ideal",
    ],
    "Five model-level differences feed the base-point discrepancy. The frozen Vbi drives the donor reversal; the interface-SRH difference is the leading suspect for the 96 mV deficit and the sweep shape mismatches; the optics difference explains the Jsc gap. Each is real, mostly by design.",
)

# ===================== SLIDE 12 =====================
figure_slide(
    "Which Absolute Is Realistic? SolarLab Sits at the Measured Median",
    "f7_experiment_voc.png",
    [
        "Shockley–Queisser ceiling at E_{g}=1.53 eV: ~1.28–1.33 V [12,14,15]",
        "Measured-device V_{oc} median: ~1.05–1.13 V; champions to ~1.17 V [16,18]",
        "SolarLab 1.072 V sits inside the measured median band [13,18]",
        "SCAPS 1.1676 V is at the published champion ceiling [16,17]",
        "SCAPS PCE 26.69% exceeds the certified MAPbI_{3} record (~21–23%) [16,17]",
        "On absolutes, SolarLab is the more conservative, representative value",
    ],
    "Neither solver is experiment, so we ask which absolute better represents a fabricated cell. SolarLab's 1.072 V implies a ~230-260 mV non-radiative deficit, typical for a routine stack. SCAPS's value implies champion-grade passivation. Tuning SolarLab up to match SCAPS would move it away from the experimental median.",
)

# ===================== SLIDE 13 =====================
text_slide(
    "Recommendations: One Cheap Fix, No V_{bi} Refactor",
    [
        "P0: re-extract all four FOM from one J–V array; quote PCE 22.1%",
        "P0: stop quoting J_{sc}=25.73 (the foreign N_{grid}=30 cell)",
        "P1: standardize swept current on charge-conserving median face",
        "Document frozen-V_{bi} as the Nd_{ETL} reversal cause; do not hot-fix",
        "AVOID: routing V_{bi,eff} to boundary or finite-S contact refactor",
        "Refactor re-baselines 17/17 guards, chases champion-ceiling parity",
    ],
    "The one near-zero-risk fix is the FOM re-extraction, which removes both the PCE artifact and the self-consistency violation. The frozen-Vbi reversal is a real defect worth documenting but not worth a multi-week boundary-condition rework that would tune SolarLab off the experimental median.",
    bsize=18,
)

# ===================== SLIDE 14 (references, two columns) =====================
refs = [
    "[1] M. Burgelman, P. Nollet, S. Degrave. Modelling polycrystalline semiconductor solar cells. Thin Solid Films 361–362, 527–532 (2000).",
    "[2] M. Burgelman, J. Marlein. Analysis of graded band gap solar cells with SCAPS. Thin Solid Films 515(15), 6276–6278 (2007).",
    "[3] K. Decock, S. Khelifi, M. Burgelman. Modelling multivalent defects in thin-film solar cells. Thin Solid Films 519(21), 7481–7484 (2011).",
    "[4] J. Niemegeers, M. Burgelman. Thermionic-emission interface transport and numerical scheme (1998).",
    "[5] D. L. Scharfetter, H. K. Gummel. Large-signal analysis of a silicon Read diode oscillator. IEEE TED 16(1), 64–77 (1969).",
    "[6] SCAPS Manual (Feb 2016), University of Gent. pp. 9–10, 24/26, 30–31, 36–37, 39–41.",
    "[7] W. Shockley, W. T. Read. Statistics of the recombinations of holes and electrons. Phys. Rev. 87(5), 835–842 (1952).",
    "[8] R. J. Pauwels, G. Vanhoutte. Influence of interface state and energy barriers on heterojunction solar cell efficiency. J. Phys. D 11(5), 649 (1978).",
    "[9] J. Nelson. The Physics of Solar Cells. Imperial College Press (2003).",
    "[10] P. Calado et al. Evidence for ion migration in hybrid perovskite solar cells. Nat. Commun. 7, 13831 (2016) / Driftfusion.",
    "[11] N. E. Courtier et al. IonMonger: a free and fast planar perovskite solar cell simulator. J. Comput. Electron. 18, 1435–1449 (2019).",
    "[12] Effect of perovskite thickness on EL and conversion efficiency. J. Phys. Chem. Lett. (2020). doi:10.1021/acs.jpclett.0c02363.",
    "[13] K. Tvingstedt et al. Radiative efficiency of lead iodide perovskite solar cells. Sci. Rep. 4, 6071 (2014). doi:10.1038/srep06071.",
    "[14] L. M. Pazos-Outón et al. Fundamental efficiency limit of lead iodide perovskite. J. Phys. Chem. Lett. 9(7), 1703–1711 (2018).",
    "[15] W. E. I. Sha et al. The efficiency limit of CH₃NH₃PbI₃ perovskite solar cells. Appl. Phys. Lett. 106, 221104 (2015). arXiv:1506.09003.",
    "[16] Mesoporous-TiO₂/spiro champion. Energy Environ. Sci. (2016); PMC4705040. J_{sc} 24.6, V_{oc} 1.16 V, FF 0.73, PCE 20.8%.",
    "[17] Low-loss MAPbI₃ stack: V_{oc} 1.21 V at 1.53 eV, certified 23.09% PCE. Sci. China Mater. (2025). doi:10.1007/s40843-025-3457-3.",
    "[18] Single-crystal MAPbI₃ cells > 21% PCE, V_{oc} 1.0–1.1 V. ACS Energy Lett. 4(5), 1258 (2019). doi:10.1021/acsenergylett.9b00847.",
]
s = new_slide()
add_title(s, "References")
mid = (len(refs) + 1) // 2
cols = [refs[:mid], refs[mid:]]
col_geom = [(0.6, 6.1), (6.85, 6.1)]
for refs_col, (cl, cw) in zip(cols, col_geom):
    tb = s.shapes.add_textbox(Inches(cl), Inches(1.55), Inches(cw), Inches(5.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, r in enumerate(refs_col):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = 1.0
        add_rich(p, r, size=11, color=SLATE)
add_notes(s, "Full numbered list of 18 entries; numbering matches the report's Section 7 exactly. The inline [n] brackets placed on slides 2, 9, 11 and 12 all resolve into this list: [1,5] solver formulation and SG flux, [2] tunnelling-off, [4] thermionic-emission interface BC, [6] SCAPS contacts/flatband, [7,8] Shockley-Read-Hall and Pauwels-Vanhoutte interface SRH, and [12-18] the experimental MAPbI3 Voc / SQ realism cluster. Literature values are used only for realism-positioning and are flagged as such in the report.")

# ---------- save + reload ----------
prs.save(OUT)
print("SAVED:", OUT)
if missing:
    print("MISSING_FIGURES:", missing)
else:
    print("MISSING_FIGURES: none")

re = Presentation(OUT)
print("RELOADED_OK slide_count =", len(re.slides.__iter__.__self__._sldIdLst))
for idx, sld in enumerate(re.slides, 1):
    print("  slide", idx, "shapes =", len(sld.shapes))
